from __future__ import annotations

from datetime import date, datetime
import io
import json
import math
from pathlib import Path
import re
import shutil
import subprocess
import tempfile
from typing import Any

from .schemas import DeckRequest
from .workbook import DetectedTable

ROOT = Path(__file__).resolve().parents[1]
GENERATOR = ROOT / "generator" / "generate_deck.cjs"

# ── AMD dark palette (used by the image renderer) ─────────────────────────────
_BG      = "#000000"
_PANEL   = "#0D0D0D"
_SEP     = "#404040"
_GRID    = "#2A2A2A"
_WHITE   = "#FFFFFF"
_BODY    = "#D9D9D9"
_AXIS    = "#A6A6A6"
_GOLD    = "#C1A968"
_SERIES  = [
    "#C1A968", "#00C2DE", "#5B8FF9", "#3CCB8C",
    "#D96AA7", "#FF6B2C", "#8F7CEC", "#7BC8A4",
    "#F4A261", "#B8A1E3",
]


def _hex_rgb(h: str):
    h = h.lstrip("#")
    from pptx.dml.color import RGBColor
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _parse_number(value: Any) -> float | None:
    try:
        return float(str(value).replace(",", "").strip())
    except (ValueError, TypeError):
        return None


def _category_sort_key(value: str) -> tuple:
    """Parse message-size categories like '1024', '8K', '1M' for sorting."""
    s = str(value).strip().lower()
    scales = {"k": 1024, "m": 1024**2, "g": 1024**3, "t": 1024**4}
    for suffix, mult in scales.items():
        for variant in (suffix + "b", suffix + "ib", suffix):
            if s.endswith(variant):
                try:
                    return (0, float(s[: -len(variant)]) * mult)
                except ValueError:
                    pass
    try:
        return (0, float(s))
    except ValueError:
        return (1, s)


def _extract_series(table: DetectedTable, slide_req) -> tuple[list[str], list[str], list[list[float | None]]]:
    """Return (sorted_categories, series_names, values_per_series)."""
    headers = table.headers
    rows    = table.rows
    cat_i   = slide_req.category_index
    ser_is  = slide_req.series_indexes or []

    buckets: dict[str, dict[str, list[float]]] = {}
    cat_order: list[str] = []

    for row in rows:
        cat = str(row[cat_i]).strip() if cat_i < len(row) else ""
        if cat not in buckets:
            buckets[cat] = {}
            cat_order.append(cat)
        for si in ser_is:
            if si >= len(row):
                continue
            v = _parse_number(row[si])
            if v is not None:
                name = headers[si] if si < len(headers) else f"Col {si}"
                buckets[cat].setdefault(name, []).append(v)

    cats = sorted(cat_order, key=_category_sort_key) if slide_req.sort_categories else cat_order
    names = [headers[si] for si in ser_is if si < len(headers)]

    def _agg(vals: list[float]) -> float | None:
        if not vals:
            return None
        return sorted(vals)[len(vals) // 2]  # median

    values = [
        [_agg(buckets.get(c, {}).get(n, [])) for c in cats]
        for n in names
    ]
    return cats, names, values


def _fmt_cat(c: str) -> str:
    key = _category_sort_key(c)
    if key[0] == 0:
        n = key[1]
        if n >= 1024**2:
            return f"{int(n/1024**2)}M"
        if n >= 1024:
            return f"{int(n/1024)}K"
        return str(int(n))
    return c


def _render_line_chart(table: DetectedTable, slide_req) -> io.BytesIO:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    cats, names, values = _extract_series(table, slide_req)
    x = np.arange(len(cats))

    fig, ax = plt.subplots(figsize=(12.8, 5.6))
    fig.patch.set_facecolor(_BG)
    ax.set_facecolor(_PANEL)

    for i, (name, vals) in enumerate(zip(names, values)):
        color = _SERIES[i % len(_SERIES)]
        y = [v if v is not None else float("nan") for v in vals]
        ax.plot(x, y, color=color, linewidth=2.2, marker="o", markersize=4.5,
                markerfacecolor=color, markeredgewidth=0, label=name, zorder=3)

    xlabels = [_fmt_cat(c) for c in cats]
    step = max(1, len(cats) // 20)
    ax.set_xticks(x[::step])
    ax.set_xticklabels(xlabels[::step], color=_BODY, fontsize=8)
    ax.tick_params(axis="x", colors=_BODY)
    ax.tick_params(axis="y", colors=_AXIS, labelsize=8)

    if slide_req.x_axis_title:
        ax.set_xlabel(slide_req.x_axis_title, color=_BODY, fontsize=10, labelpad=5)
    ax.set_ylabel(slide_req.y_axis_title or "Value", color=_BODY, fontsize=10, labelpad=5)

    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    ax.spines["bottom"].set_color(_SEP)
    ax.spines["left"].set_color(_SEP)
    ax.grid(axis="y", color=_GRID, linewidth=0.7, zorder=0)
    ax.grid(axis="x", color=_GRID, linewidth=0.4, linestyle="--", zorder=0)
    ax.set_xlim(-0.5, max(len(cats) - 0.5, 0.5))

    ax.legend(loc="upper left", frameon=True, fontsize=8,
              facecolor="#111111", edgecolor=_SEP, labelcolor=_BODY,
              ncol=min(len(names), 4), handlelength=1.4, columnspacing=0.9)

    fig.tight_layout(pad=1.0)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, facecolor=_BG, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _render_grouped_bar(table: DetectedTable, slide_req) -> io.BytesIO:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    cats, names, values = _extract_series(table, slide_req)
    x = np.arange(len(cats))
    n = len(names)
    width = min(0.8 / max(n, 1), 0.25)

    fig, ax = plt.subplots(figsize=(12.8, 5.6))
    fig.patch.set_facecolor(_BG)
    ax.set_facecolor(_PANEL)

    for i, (name, vals) in enumerate(zip(names, values)):
        offset = (i - n / 2 + 0.5) * width
        color  = _SERIES[i % len(_SERIES)]
        y = [v if v is not None else 0 for v in vals]
        ax.bar(x + offset, y, width=width * 0.88, color=color, label=name, zorder=3)

    ax.set_xticks(x)
    ax.set_xticklabels(cats, color=_BODY, fontsize=8,
                       rotation=30 if len(cats) > 8 else 0, ha="right")
    ax.tick_params(axis="x", colors=_BODY)
    ax.tick_params(axis="y", colors=_AXIS, labelsize=8)
    ax.set_ylabel(slide_req.y_axis_title or "Value", color=_BODY, fontsize=10, labelpad=5)

    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    ax.spines["bottom"].set_color(_SEP)
    ax.spines["left"].set_color(_SEP)
    ax.grid(axis="y", color=_GRID, linewidth=0.7, zorder=0)
    ax.legend(loc="best", frameon=True, fontsize=8, facecolor="#111111",
              edgecolor=_SEP, labelcolor=_BODY, ncol=min(n, 4), handlelength=1.4)

    fig.tight_layout(pad=1.0)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, facecolor=_BG, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _render_data_matrix(table: DetectedTable, slide_req) -> io.BytesIO:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    headers = table.headers[:12]
    rows    = [r[:12] for r in table.rows[:20]]

    fig, ax = plt.subplots(figsize=(12.8, max(2.5, len(rows) * 0.35 + 0.8)))
    fig.patch.set_facecolor(_BG)
    ax.axis("off")

    cell_text = [[str(v) if v is not None else "" for v in row] for row in rows]
    tbl = ax.table(cellText=cell_text, colLabels=headers, cellLoc="center", loc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(8)
    tbl.scale(1, 1.8)

    for (row, col), cell in tbl.get_celld().items():
        cell.set_edgecolor("#1F1F1F")
        if row == 0:
            cell.set_facecolor("#151515")
            cell.get_text().set_color(_GOLD)
            cell.get_text().set_fontweight("bold")
        else:
            cell.set_facecolor(_BG if row % 2 == 0 else _PANEL)
            cell.get_text().set_color(_BODY)

    fig.tight_layout(pad=0.4)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, facecolor=_BG, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_pptx_slide(prs, title_text: str, chart_buf: io.BytesIO) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _hex_rgb(_BG)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.26), Inches(11.85), Inches(0.60))
    tf = tb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    font_pt = 17 if len(title_text) > 95 else (19 if len(title_text) > 65 else 22)
    run.font.size      = Pt(font_pt)
    run.font.bold      = True
    run.font.color.rgb = _hex_rgb(_WHITE)
    run.font.name      = "Arial"

    # Separator line
    sep = slide.shapes.add_shape(1, Inches(0.55), Inches(0.95), Inches(12.05), Pt(0.65))
    sep.fill.solid()
    sep.fill.fore_color.rgb = _hex_rgb(_SEP)
    sep.line.fill.background()

    # Chart image
    top = Inches(1.06)
    slide.shapes.add_picture(
        chart_buf,
        Inches(0.35), top,
        slide_w - Inches(0.7), slide_h - top - Inches(0.08),
    )


def generate_deck_images(deck_request: DeckRequest, tables: dict[str, DetectedTable]) -> tuple[bytes, str]:
    """
    Image-based renderer: charts are drawn by matplotlib and embedded as PNGs.
    Produces a .pptx that opens in all PowerPoint versions (no chart XML objects).
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _renderers = {
        "line":        _render_line_chart,
        "grouped_bar": _render_grouped_bar,
        "difference":  _render_grouped_bar,   # fallback: bar style for diffs
        "data_matrix": _render_data_matrix,
    }

    for slide_req in deck_request.slides:
        table = tables.get(slide_req.table_id)
        if table is None:
            raise ValueError(f"The source table '{slide_req.table_id}' is no longer available.")
        renderer = _renderers.get(slide_req.slide_type, _render_line_chart)
        chart_buf = renderer(table, slide_req)
        _add_pptx_slide(prs, slide_req.title, chart_buf)

    output_filename = _safe_filename(deck_request.output_filename)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), output_filename


def _safe_filename(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", value).strip("._")
    if not cleaned:
        cleaned = "presentation"
    if not cleaned.lower().endswith(".pptx"):
        cleaned += ".pptx"
    return cleaned


def _json_value(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, (date, datetime)):
        return value.isoformat()
    return str(value)


def _table_payload(table: DetectedTable) -> dict[str, Any]:
    return {
        "table_id": table.table_id,
        "sheet_name": table.sheet_name,
        "kind": table.kind,
        "title": table.title,
        "display_name": table.display_name,
        "header_row": table.header_row,
        "data_start_row": table.data_start_row,
        "data_end_row": table.data_end_row,
        "start_col": table.start_col,
        "end_col": table.end_col,
        "headers": list(table.headers),
        "rows": [[_json_value(value) for value in row] for row in table.rows],
        "column_profiles": [profile.public_dict() for profile in table.column_profiles],
    }


def _find_node() -> str:
    executable = shutil.which("node") or shutil.which("node.exe")
    if not executable:
        raise RuntimeError(
            "Node.js was not found. Install Node.js 18 or newer, then restart the application. "
            "No npm install is required because the PowerPoint renderer is bundled."
        )
    return executable


def generate_deck(deck_request: DeckRequest, tables: dict[str, DetectedTable]) -> tuple[bytes, str]:
    if not GENERATOR.exists():
        raise RuntimeError(f"The bundled PowerPoint generator is missing: {GENERATOR}")

    selected_tables: dict[str, dict[str, Any]] = {}
    for slide in deck_request.slides:
        table = tables.get(slide.table_id)
        if table is None:
            raise ValueError(f"The source table '{slide.table_id}' is no longer available.")
        selected_tables[slide.table_id] = _table_payload(table)

    payload = {
        "deck_title": deck_request.deck_title,
        "slides": [slide.model_dump(mode="json") for slide in deck_request.slides],
        "tables": selected_tables,
    }
    output_filename = _safe_filename(deck_request.output_filename)

    with tempfile.TemporaryDirectory(prefix="slide_generator_") as temporary_directory:
        work = Path(temporary_directory)
        input_path = work / "presentation_plan.json"
        output_path = work / output_filename
        input_path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")

        process = subprocess.run(
            [_find_node(), str(GENERATOR), str(input_path), str(output_path)],
            cwd=ROOT,
            capture_output=True,
            text=True,
            timeout=240,
            check=False,
        )
        if process.returncode != 0:
            details = (process.stderr or process.stdout or "Unknown renderer error").strip()
            raise RuntimeError(details[-6000:])
        if not output_path.exists() or output_path.stat().st_size == 0:
            raise RuntimeError("The PowerPoint renderer completed without producing an output file.")
        return output_path.read_bytes(), output_filename
